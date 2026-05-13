import os
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
import openpyxl
from pptx import Presentation

# i created this to handle reading text from files
# the problem was each file type is completely different
# pdf needs one library docx needs another images need ocr
# so i made separate functions for each and one main router

# tesseract path - needed this because windows doesnt
# find it automatically without setting the path
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# file types my system currently supports
# i can add more here later without changing anything else
SUPPORTED_TYPES = ['.pdf', '.docx', '.jpg', 
                   '.jpeg', '.png', '.txt', 
                   '.xlsx', '.pptx']


def get_file_type(file_path):
    # splitext gives me the extension part of filename
    # lower() so .PDF and .pdf are treated the same
    return os.path.splitext(file_path)[1].lower()


def is_supported(file_path):
    return get_file_type(file_path) in SUPPORTED_TYPES


def extract_from_pdf(file_path):
    # i tried PyPDF2 first but it missed text sometimes
    # switched to pdfplumber which is more accurate
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + " "
    except Exception as e:
        print(f"couldnt read pdf: {e}")
        return None
    return text.strip()


def extract_from_docx(file_path):
    # word documents are structured as paragraphs
    # so i loop through each one and collect the text
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + " "
    except Exception as e:
        print(f"couldnt read docx: {e}")
        return None
    return text.strip()


def extract_from_image(file_path):
    # this was new to me - OCR reads text from images
    # tesseract is the engine behind it made by google
    # it looks at the image pixels and figures out letters
    text = ""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
    except Exception as e:
        print(f"couldnt read image: {e}")
        return None
    return text.strip()


def extract_from_txt(file_path):
    # simplest one - just open and read
    # utf-8 handles special characters
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        print(f"couldnt read txt: {e}")
        return None


def extract_from_xlsx(file_path):
    # excel has sheets inside sheets has rows
    # rows have cells - i go through all three levels
    # str() converts numbers to text so nothing is missed
    text = ""
    try:
        workbook = openpyxl.load_workbook(file_path)
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + " "
    except Exception as e:
        print(f"couldnt read xlsx: {e}")
        return None
    return text.strip()


def extract_from_pptx(file_path):
    # slides have shapes - text boxes images charts etc
    # not all shapes have text so i check with hasattr
    # otherwise it crashes on image shapes
    text = ""
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text += shape.text + " "
    except Exception as e:
        print(f"couldnt read pptx: {e}")
        return None
    return text.strip()


def extract_text(file_path):
    # main function - everything in pipeline calls this
    # it figures out file type and calls right function
    # returns text if successful or None if it failed

    print(f"reading: {os.path.basename(file_path)}")

    if not is_supported(file_path):
        print(f"file type not supported")
        return None

    extension = get_file_type(file_path)

    if extension == ".pdf":
        text = extract_from_pdf(file_path)
    elif extension == ".docx":
        text = extract_from_docx(file_path)
    elif extension in [".jpg", ".jpeg", ".png"]:
        text = extract_from_image(file_path)
    elif extension == ".txt":
        text = extract_from_txt(file_path)
    elif extension == ".xlsx":
        text = extract_from_xlsx(file_path)
    elif extension == ".pptx":
        text = extract_from_pptx(file_path)
    else:
        return None

    # if text is too short something probably went wrong
    if not text or len(text.strip()) < 10:
        print("couldnt find readable text")
        return None

    print(f"got {len(text)} characters")
    return text


if __name__ == "__main__":

    print("testing extractor\n")

    test_files = os.listdir("uploads")

    if len(test_files) == 0:
        print("uploads folder is empty")
        print("add a file to test")
    else:
        for file in test_files:
            file_path = os.path.join("uploads", file)

            if os.path.isdir(file_path):
                continue

            print(f"testing: {file}")
            print("-" * 40)

            text = extract_text(file_path)

            if text:
                print(f"\nfirst 300 characters:")
                print(text[:300])
            else:
                print("extraction failed")
            print()